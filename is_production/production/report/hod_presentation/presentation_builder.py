from __future__ import annotations

from io import BytesIO
from typing import BinaryIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5
FONT_NAME = "Aptos"

BLACK = "050505"
CARD = "151515"
CARD_ALT = "1D1D1D"
CARD_BORDER = "3A3A3A"
NAVY = "0F1F53"
GREEN = "18A957"
RED = "E03124"
ORANGE = "F59E0B"
UTIL_GRAY = "737373"
WHITE = "FFFFFF"
MUTED = "AAB4C3"
LIGHT_TEXT = "E6EAF0"
GRID = "353535"
AVAIL_TARGET_LINE = "FF2D2D"
UTIL_TARGET_LINE = "7AC943"
SPARE_PURPLE = "D291FF"


def build_hod_presentation(
    payload: dict,
    output: BinaryIO | BytesIO,
) -> None:
    """Build a black HOD presentation from the exact selected report payload."""

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    site_payloads = payload.get("site_payloads") or [payload]
    site_payloads = [row for row in site_payloads if isinstance(row, dict)]

    if not site_payloads:
        site_payloads = [payload]

    combined_site_name = payload.get("site") or " / ".join(
        str(row.get("site") or "") for row in site_payloads
    )

    first_payload = site_payloads[0]
    title_payload = dict(first_payload)
    title_payload["site"] = combined_site_name
    title_payload["period_label"] = (
        payload.get("period_label")
        or first_payload.get("period_label", "")
    )
    title_payload["generated_by"] = (
        payload.get("generated_by")
        or first_payload.get("generated_by", "")
    )
    title_payload["generated_at"] = (
        payload.get("generated_at")
        or first_payload.get("generated_at", "")
    )

    prs.core_properties.title = f"HOD Presentation - {combined_site_name}"
    prs.core_properties.subject = (
        "Filtered HOD production, excavator, availability and utilisation report"
    )
    prs.core_properties.author = title_payload.get("generated_by") or "Isambane"

    _add_title_slide(prs, title_payload, len(site_payloads))

    summary_chunks = [
        site_payloads[index : index + 3]
        for index in range(0, len(site_payloads), 3)
    ]

    for index, chunk in enumerate(summary_chunks, start=1):
        _add_production_summary_slide(
            prs,
            chunk,
            title_payload,
            part=index,
            total_parts=len(summary_chunks),
        )

    for site_payload in site_payloads:
        _add_site_au_slides(prs, site_payload)

    _apply_footer_numbers(prs)
    prs.save(output)


def _add_title_slide(prs: Presentation, payload: dict, site_count: int) -> None:
    slide = _blank_slide(prs)

    _add_rect(slide, 0.0, 0.0, 0.18, SLIDE_HEIGHT, GREEN, GREEN)
    _add_rect(slide, 0.72, 1.00, 1.30, 0.08, GREEN, GREEN)

    _add_text(
        slide,
        "HOD PRODUCTION PRESENTATION",
        0.72,
        1.28,
        11.90,
        0.62,
        30,
        WHITE,
        bold=True,
    )
    _add_text(
        slide,
        payload.get("site", ""),
        0.72,
        2.05,
        11.90,
        0.62,
        25,
        LIGHT_TEXT,
        bold=True,
    )
    _add_text(
        slide,
        "Production Summary and Availability & Utilisation",
        0.72,
        2.78,
        11.90,
        0.38,
        16,
        MUTED,
    )

    _add_rect(slide, 0.72, 3.65, 5.66, 1.23, CARD, CARD_BORDER, radius=True)
    _add_text(slide, "REPORTING PERIOD", 1.00, 3.88, 2.30, 0.24, 9, MUTED, bold=True)
    _add_text(
        slide,
        payload.get("period_label", ""),
        1.00,
        4.18,
        5.00,
        0.35,
        17,
        WHITE,
        bold=True,
    )

    availability = payload.get("availability", {}) or {}
    _add_rect(slide, 6.70, 3.65, 5.90, 1.70, CARD, CARD_BORDER, radius=True)
    _add_filter_line(slide, "SITES", str(site_count), 7.00, 3.86)
    _add_filter_line(
        slide,
        "A&U VIEW",
        availability.get("summary_type", ""),
        7.00,
        4.26,
    )
    _add_filter_line(
        slide,
        "MACHINES",
        availability.get("machine_scope", ""),
        7.00,
        4.66,
    )
    _add_filter_line(
        slide,
        "TARGET",
        availability.get("au_target_filter", ""),
        7.00,
        5.06,
    )

    generated_line = "Generated {0} by {1}".format(
        payload.get("generated_at", ""),
        payload.get("generated_by", ""),
    )
    _add_text(slide, generated_line, 0.72, 6.86, 11.90, 0.22, 8, MUTED)


def _add_production_summary_slide(
    prs: Presentation,
    site_payloads: list[dict],
    title_payload: dict,
    *,
    part: int,
    total_parts: int,
) -> None:
    slide = _blank_slide(prs)
    title = "HOD PRODUCTION SUMMARY"
    if total_parts > 1:
        title += f" ({part}/{total_parts})"

    _add_section_header(
        slide,
        title,
        title_payload.get("period_label", ""),
        f"{len(site_payloads)} Site{'s' if len(site_payloads) != 1 else ''}",
    )

    count = len(site_payloads)
    gap = 0.16
    left = 0.18
    total_width = 12.97
    card_width = (total_width - gap * (count - 1)) / count

    for index, site_payload in enumerate(site_payloads):
        x = left + index * (card_width + gap)
        _add_site_summary_card(slide, site_payload, x, 1.00, card_width, 6.15)


def _add_site_summary_card(
    slide,
    payload: dict,
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    production = payload.get("production", {}) or {}
    excavators = payload.get("excavators", {}) or {}

    forecast_variance = _num(production.get("forecast_variance_bcm"))
    variance_colour = GREEN if forecast_variance >= 0 else RED
    variance_fill = "123623" if forecast_variance >= 0 else "411919"

    _add_rect(slide, x, y, w, h, CARD, CARD_BORDER, radius=True)
    _add_rect(slide, x + 0.08, y + 0.08, w - 0.16, 0.38, NAVY, NAVY, radius=True)
    _add_text(
        slide,
        payload.get("site", ""),
        x + 0.15,
        y + 0.12,
        w - 0.30,
        0.26,
        12,
        WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    _add_rect(
        slide,
        x + 0.08,
        y + 0.54,
        w - 0.16,
        0.60,
        variance_fill,
        variance_fill,
        radius=True,
    )
    _add_text(
        slide,
        "FORECAST VARIANCE",
        x + 0.16,
        y + 0.63,
        w - 0.32,
        0.16,
        7,
        MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        _fmt_signed(forecast_variance, 0, " BCM"),
        x + 0.16,
        y + 0.80,
        w - 0.32,
        0.26,
        17,
        variance_colour,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    total_hours = _num(excavators.get("total_hours"))
    non_prod_hours = _num(excavators.get("non_production_hours"))
    production_hours = _num(excavators.get("production_hours"))
    actual_bcm = _num(production.get("actual_bcm"))

    metric_rows = [
        ("Monthly target", "BCM", _fmt_number(production.get("monthly_target_bcm"), 0), WHITE),
        ("Forecast", "BCM", _fmt_number(production.get("forecast_bcm"), 0), WHITE),
        ("Waste variance", "BCM", _fmt_signed(production.get("waste_variance_bcm"), 0), _variance_colour(production.get("waste_variance_bcm"))),
        ("Coal variance", "TONS", _fmt_signed(production.get("coal_variance_tons"), 0), _variance_colour(production.get("coal_variance_tons"))),
        ("Actual BCMs", "BCM", _fmt_number(actual_bcm, 0), WHITE),
        ("Actual coal", "TONS", _fmt_number(production.get("actual_coal_tons"), 0), WHITE),
        ("Daily required", "BCM", _fmt_number(production.get("daily_required_bcm"), 1), WHITE),
        ("Daily achieved", "BCM", _fmt_number(production.get("daily_achieved_bcm"), 1), WHITE),
        ("Excavator hours (From Pre-Use)", "HRS", _fmt_number(total_hours, 1), WHITE),
        ("Non-production hours", "HRS", _fmt_number(non_prod_hours, 1), ORANGE if non_prod_hours > 0 else MUTED),
        ("Excavator hours (Production Hours)", "HRS", _fmt_number(production_hours, 1), WHITE),
        (
            "Average BCM/H ({0} BCM / {1} HRS)".format(
                _fmt_number(actual_bcm, 0),
                _fmt_number(production_hours, 1),
            ),
            "",
            _fmt_number(payload.get("average_bcm_h"), 1),
            GREEN if _num(payload.get("average_bcm_h")) > 0 else RED,
        ),
        (
            "Days worked / left",
            "",
            "{0} / {1}".format(
                _fmt_number(production.get("days_worked"), 0),
                _fmt_number(production.get("days_left"), 0),
            ),
            WHITE,
        ),
        ("Strip ratio", "", _fmt_number(production.get("strip_ratio"), 1), WHITE),
    ]

    table_y = y + 1.25
    row_height = 0.302
    label_w = w * 0.58
    unit_w = w * 0.14
    value_w = w - label_w - unit_w - 0.16

    for row_index, (label, unit, value, value_colour) in enumerate(metric_rows):
        row_y = table_y + row_index * row_height
        fill = CARD_ALT if row_index % 2 else CARD
        _add_rect(slide, x + 0.08, row_y, w - 0.16, row_height, fill, CARD_BORDER)
        _add_text(slide, label, x + 0.14, row_y + 0.02, label_w - 0.08, row_height - 0.04, 6.8, LIGHT_TEXT)
        _add_text(
            slide,
            unit,
            x + 0.08 + label_w,
            row_y + 0.02,
            unit_w - 0.02,
            row_height - 0.04,
            6.8,
            MUTED,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
        _add_text(
            slide,
            value,
            x + 0.08 + label_w + unit_w,
            row_y + 0.02,
            value_w,
            row_height - 0.04,
            7.2,
            value_colour,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )

    forecast_delivery = _num(production.get("forecast_delivery_percent"))
    delivery_colour = GREEN if forecast_delivery >= 100 else ORANGE
    _add_text(
        slide,
        "Forecast delivery",
        x + 0.14,
        y + h - 0.38,
        1.55,
        0.20,
        7,
        MUTED,
    )
    _add_text(
        slide,
        _fmt_number(forecast_delivery, 1) + "%",
        x + 1.65,
        y + h - 0.38,
        0.90,
        0.20,
        7.5,
        delivery_colour,
        bold=True,
    )


def _add_site_au_slides(prs: Presentation, payload: dict) -> None:
    availability = payload.get("availability", {}) or {}
    machine_series = availability.get("machine_series", {}) or {}
    adt_rows = machine_series.get("ADT") or []

    chunks = [adt_rows[index : index + 16] for index in range(0, len(adt_rows), 16)]
    if not chunks:
        chunks = [[]]

    for part, chunk in enumerate(chunks, start=1):
        _add_site_au_slide(
            prs,
            payload,
            chunk,
            part=part,
            total_parts=len(chunks),
        )


def _add_site_au_slide(
    prs: Presentation,
    payload: dict,
    adt_rows: list[dict],
    *,
    part: int,
    total_parts: int,
) -> None:
    slide = _blank_slide(prs)
    site = payload.get("site", "")
    title = "AVAILABILITY & UTILISATION"
    if total_parts > 1:
        title += f" ({part}/{total_parts})"

    _add_section_header(
        slide,
        title,
        payload.get("period_label", ""),
        site,
    )

    availability = payload.get("availability", {}) or {}
    filter_text = "{0}  |  {1}  |  {2}".format(
        availability.get("summary_type", ""),
        availability.get("machine_scope", ""),
        availability.get("au_target_filter", ""),
    )
    _add_text(slide, filter_text, 0.30, 0.82, 12.70, 0.20, 8, MUTED, bold=True)

    categories = availability.get("categories", []) or []
    for index, category in enumerate(categories[:10]):
        col = index % 5
        row = index // 5
        x = 0.26 + col * 2.57
        y = 1.10 + row * 0.83
        _add_category_metric_card(slide, category, x, y, 2.40, 0.70, availability)

    _add_adt_chart(
        slide,
        adt_rows,
        availability_target=_num(availability.get("availability_target")) or 85,
        utilisation_target=_num(availability.get("utilisation_target")) or 80,
    )


def _add_category_metric_card(
    slide,
    category: dict,
    x: float,
    y: float,
    w: float,
    h: float,
    availability: dict,
) -> None:
    _add_rect(slide, x, y, w, h, CARD, CARD_BORDER, radius=True)
    _add_text(
        slide,
        category.get("title", ""),
        x + 0.08,
        y + 0.06,
        w - 0.16,
        0.16,
        7,
        WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    _add_metric_bubble(
        slide,
        x + 0.28,
        y + 0.26,
        0.36,
        "A",
        category.get("availability"),
        availability.get("availability_target", 85),
    )
    _add_metric_bubble(
        slide,
        x + w - 0.64,
        y + 0.26,
        0.36,
        "U",
        category.get("utilisation"),
        availability.get("utilisation_target", 80),
    )

    _add_text(
        slide,
        _fmt_percent_or_nd(category.get("availability")),
        x + 0.68,
        y + 0.30,
        0.53,
        0.20,
        7.2,
        _status_colour(category.get("availability"), availability.get("availability_target", 85)),
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        _fmt_percent_or_nd(category.get("utilisation")),
        x + 1.20,
        y + 0.30,
        0.53,
        0.20,
        7.2,
        _status_colour(category.get("utilisation"), availability.get("utilisation_target", 80)),
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def _add_metric_bubble(
    slide,
    x: float,
    y: float,
    size: float,
    label: str,
    value,
    target,
) -> None:
    colour = _status_colour(value, target)
    fill = "123623" if colour == GREEN else "411919"
    if value is None:
        fill = CARD_ALT
        colour = MUTED

    _add_oval(slide, x, y, size, size, fill, colour)
    _add_text(
        slide,
        label,
        x,
        y + 0.06,
        size,
        0.15,
        6.5,
        colour,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def _add_adt_chart(
    slide,
    rows: list[dict],
    *,
    availability_target: float,
    utilisation_target: float,
) -> None:
    _add_rect(slide, 0.25, 2.78, 12.83, 4.22, CARD, CARD_BORDER, radius=True)
    _add_text(
        slide,
        "ADT AVAILABILITY & UTILISATION - SELECTED FILTERS",
        0.55,
        2.90,
        8.50,
        0.24,
        11,
        WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    _add_rect(slide, 9.30, 2.94, 0.18, 0.12, ORANGE, ORANGE)
    _add_text(slide, "Availability", 9.55, 2.89, 1.10, 0.20, 7, LIGHT_TEXT)
    _add_rect(slide, 10.78, 2.94, 0.18, 0.12, UTIL_GRAY, UTIL_GRAY)
    _add_text(slide, "Utilisation", 11.03, 2.89, 1.05, 0.20, 7, LIGHT_TEXT)

    plot_x = 0.85
    plot_y = 3.35
    plot_w = 11.75
    plot_h = 2.95
    base_y = plot_y + plot_h

    for tick in range(0, 101, 20):
        y = plot_y + plot_h - (tick / 100.0) * plot_h
        _add_rect(slide, plot_x, y, plot_w, 0.008, GRID, GRID)
        _add_text(
            slide,
            f"{tick}%",
            0.33,
            y - 0.10,
            0.43,
            0.20,
            6.5,
            MUTED,
            align=PP_ALIGN.RIGHT,
        )

    avail_y = plot_y + plot_h - (availability_target / 100.0) * plot_h
    util_y = plot_y + plot_h - (utilisation_target / 100.0) * plot_h
    _add_rect(slide, plot_x, avail_y, plot_w, 0.022, AVAIL_TARGET_LINE, AVAIL_TARGET_LINE)
    _add_rect(slide, plot_x, util_y, plot_w, 0.022, UTIL_TARGET_LINE, UTIL_TARGET_LINE)

    if not rows:
        _add_text(
            slide,
            "No ADT availability and utilisation data found for the selected filters.",
            plot_x,
            4.40,
            plot_w,
            0.40,
            14,
            MUTED,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        return

    count = len(rows)
    group_w = plot_w / max(count, 1)
    bar_w = min(0.25, max(0.07, group_w * 0.27))
    gap = min(0.07, max(0.02, group_w * 0.08))
    label_size = 7 if count <= 12 else 5.8

    for index, row in enumerate(rows):
        group_left = plot_x + index * group_w
        center_x = group_left + group_w / 2
        avail_x = center_x - bar_w - gap / 2
        util_x = center_x + gap / 2

        avail = _percent(row.get("availability"))
        util = _percent(row.get("utilisation"))
        avail_h = plot_h * ((avail or 0) / 100.0)
        util_h = plot_h * ((util or 0) / 100.0)

        if avail is not None:
            _add_rect(slide, avail_x, base_y - max(avail_h, 0.02), bar_w, max(avail_h, 0.02), ORANGE, ORANGE)
        if util is not None:
            _add_rect(slide, util_x, base_y - max(util_h, 0.02), bar_w, max(util_h, 0.02), UTIL_GRAY, UTIL_GRAY)

        label_colour = SPARE_PURPLE if row.get("is_spare") else LIGHT_TEXT
        _add_text(
            slide,
            _truncate_label(row.get("machine", ""), 11),
            group_left,
            base_y + 0.08,
            group_w,
            0.30,
            label_size,
            label_colour,
            bold=True,
            align=PP_ALIGN.CENTER,
        )


def _add_section_header(slide, title: str, subtitle: str, badge: str) -> None:
    _add_rect(slide, 0.16, 0.12, 13.00, 0.66, NAVY, NAVY, radius=True)
    _add_text(slide, title, 0.42, 0.21, 8.60, 0.26, 16, WHITE, bold=True)
    _add_text(slide, subtitle, 0.42, 0.48, 8.60, 0.16, 7, LIGHT_TEXT)
    _add_rect(slide, 11.65, 0.28, 1.16, 0.28, WHITE, WHITE, radius=True)
    _add_text(
        slide,
        badge,
        11.72,
        0.31,
        1.02,
        0.18,
        7,
        NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def _add_filter_line(slide, label: str, value: str, x: float, y: float) -> None:
    _add_text(slide, label, x, y, 1.15, 0.18, 7, MUTED, bold=True)
    _add_text(slide, value, x + 1.25, y - 0.01, 4.20, 0.20, 9, WHITE, bold=True)


def _apply_footer_numbers(prs: Presentation) -> None:
    total = len(prs.slides)
    for index, slide in enumerate(prs.slides, start=1):
        if index == 1:
            continue
        _add_text(
            slide,
            f"HOD Presentation  |  {index} of {total}",
            10.10,
            7.18,
            2.70,
            0.16,
            6.5,
            MUTED,
            align=PP_ALIGN.RIGHT,
        )


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(BLACK)
    return slide


def _add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill_colour: str,
    line_colour: str,
    *,
    radius: bool = False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_colour)
    shape.line.color.rgb = _rgb(line_colour)
    shape.line.width = Pt(0.7)
    return shape


def _add_oval(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill_colour: str,
    line_colour: str,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_colour)
    shape.line.color.rgb = _rgb(line_colour)
    shape.line.width = Pt(0.9)
    return shape


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: float,
    colour: str,
    *,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)

    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)

    run = paragraph.add_run()
    run.text = str(text or "")
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(colour)
    return box


def _variance_colour(value) -> str:
    return GREEN if _num(value) >= 0 else RED


def _status_colour(value, target) -> str:
    if value is None:
        return MUTED
    return GREEN if _num(value) >= _num(target) else RED


def _fmt_number(value, decimals: int = 0) -> str:
    return f"{_num(value):,.{decimals}f}"


def _fmt_percent_or_nd(value) -> str:
    if value is None:
        return "N/D"
    return f"{_num(value):.1f}%"


def _fmt_signed(value, decimals: int = 0, suffix: str = "") -> str:
    number = _num(value)
    sign = "+" if number >= 0 else "-"
    return f"{sign}{abs(number):,.{decimals}f}{suffix}"


def _percent(value):
    if value is None:
        return None
    return max(0.0, min(100.0, _num(value)))


def _num(value) -> float:
    try:
        return float(value or 0)
    except (TypeError, ValueError):
        return 0.0


def _truncate_label(value, maximum: int) -> str:
    text = str(value or "")
    if len(text) <= maximum:
        return text
    return text[: max(1, maximum - 1)] + "..."


def _rgb(hex_colour: str) -> RGBColor:
    text = str(hex_colour or "000000").strip().lstrip("#")
    return RGBColor(
        int(text[0:2], 16),
        int(text[2:4], 16),
        int(text[4:6], 16),
    )